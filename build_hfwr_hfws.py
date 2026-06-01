"""
Build HFWR and HFWS output folders.
Each song gets its own subfolder containing:
  - SongName.pptx  (full slide deck)
  - SongName - Print.pdf  (print sheet)
Also generates hfwr_hfws_catalogue.json for the web app.
"""

import re, shutil, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BASE    = Path(r'C:\Users\noahr\OneDrive\Documents\Cowork OS\castleberry-hymns\hymn-slides')
ROOT    = BASE.parent

def slide_sort_key(p):
    stem = p.stem.upper()
    if stem.endswith('TITLE'):           return (0, 0, 0)
    m = re.search(r'VERSE\s+(\d+)\.(\d+)', stem)
    if m: return (1, int(m.group(1)), int(m.group(2)))
    m = re.search(r'CHORUS[.\s]+(\d+)', stem)
    if m: return (2, 0, int(m.group(1)))
    m = re.search(r'\.(\d+)$', stem)
    if m: return (3, 0, int(m.group(1)))
    return (9, 0, 0)

def clean_name(stem, prefix):
    name = re.sub(rf'^{re.escape(prefix)}\s*-\s*', '', stem, flags=re.IGNORECASE)
    name = re.sub(r'\s*-\s*(TITLE|VERSE.*|CHORUS.*|CODA.*|BRIDGE.*|REFRAIN.*)$',
                  '', name, flags=re.IGNORECASE).strip()
    name = name.title()
    name = re.sub(r"'S\b", "'s", name)
    return name

def to_url(path):
    return str(path.relative_to(ROOT)).replace('\\', '/')

catalogue = []

for lib in ['HFWR', 'HFWS']:
    src  = BASE / lib
    out  = BASE / f'{lib} Output'
    out.mkdir(exist_ok=True)
    dirs = sorted([d for d in src.iterdir() if d.is_dir()])
    print(f'\nProcessing {lib}: {len(dirs)} songs...')

    for i, folder in enumerate(dirs, 1):
        pngs = sorted(
            [p for p in folder.glob('*.png') if 'PRINT' not in p.stem.upper()],
            key=slide_sort_key
        )
        if not pngs:
            continue

        name     = clean_name(pngs[0].stem, folder.name)
        song_dir = out / name
        song_dir.mkdir(exist_ok=True)

        # Build PPTX (skip if already exists)
        pptx_path = song_dir / f'{name}.pptx'
        if not pptx_path.exists():
            prs              = Presentation()
            prs.slide_width  = SLIDE_W
            prs.slide_height = SLIDE_H
            blank            = prs.slide_layouts[6]
            for png in pngs:
                sl = prs.slides.add_slide(blank)
                sl.shapes.add_picture(str(png), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
            prs.save(pptx_path)

        # Copy print PDF (prefer LARGE)
        pdf_url = None
        for pdf in sorted(folder.glob('*.PDF')):
            if 'LARGE' in pdf.stem.upper():
                dest = song_dir / f'{name} - Print.pdf'
                if not dest.exists():
                    shutil.copy2(pdf, dest)
                pdf_url = to_url(dest)
                break

        catalogue.append({
            'name'  : name,
            'number': folder.name,
            'source': lib,
            'pptx'  : to_url(pptx_path),
            'pdf'   : pdf_url,
        })

        if i % 100 == 0:
            print(f'  {i}/{len(dirs)} done...')

    print(f'  {lib} complete — {sum(1 for e in catalogue if e["source"] == lib)} songs.')

cat_path = ROOT / 'hfwr_hfws_catalogue.json'
with open(cat_path, 'w', encoding='utf-8') as f:
    json.dump(catalogue, f, indent=2, ensure_ascii=False)

print(f'\nDone — {len(catalogue)} songs saved to hfwr_hfws_catalogue.json')
