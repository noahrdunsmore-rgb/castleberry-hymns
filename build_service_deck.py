"""
Castleberry Church of Christ — Worship Service Deck Builder
============================================================
Assembles a complete Sunday-service PowerPoint from a *lineup spec* (JSON),
following the worship leader's service order.

A lineup is an ordered list of items. Each item is one of:

  • Song      — pulled from an already-built per-song deck (HFWR / HFWS /
                eChoice / hand-made).  HFWR & HFWS songs support per-verse
                selection (e.g. "verse 1 only") because their source PNGs are
                on disk.
  • Frame     — a recurring liturgical slide (Call to Worship, Communion,
                Giving, Prayer, …) rendered from service-template/template.json.
  • Scripture — a reading rendered from a reference + text supplied in the
                lineup.

Every source slide in this project is a single full-bleed picture (verified),
so merging is done by copying each picture onto a fresh blank slide. That keeps
the output identical to the source decks and avoids fragile XML cloning.

Usage:
    python build_service_deck.py lineups/example-service.json
    python build_service_deck.py lineups/example-service.json -o "service-decks/CB 11 (6.7.26).pptx"

Run with --list-songs to dump every resolvable song title.
"""

from __future__ import annotations

import argparse
import io
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT          = Path(__file__).parent
SLIDES_DIR    = ROOT / "hymn-slides"
TEMPLATE_PATH = ROOT / "service-template" / "template.json"
OUT_DIR       = ROOT / "service-decks"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PICTURE_SHAPE = 13  # MSO_SHAPE_TYPE.PICTURE


# ── Song index ────────────────────────────────────────────────────────────────

def normalize(name: str) -> str:
    """Lowercase, strip punctuation/whitespace — for fuzzy title matching."""
    return re.sub(r"[^a-z0-9]", "", name.lower())


def build_song_index() -> dict:
    """
    Map normalized-title -> list of song entries. Each entry:
        {name, number, source, pptx, raw_dir}
    raw_dir is set only for HFWR/HFWS (where per-verse PNGs exist).
    Multiple entries can share a title (e.g. same hymn in two hymnals); the
    lineup can disambiguate with "source"/"number".
    """
    index: dict[str, list[dict]] = {}

    def add(entry: dict):
        index.setdefault(normalize(entry["name"]), []).append(entry)

    # 1) HFWR / HFWS Output (from the built catalogue) — verse-selectable.
    cat_path = ROOT / "hymn_catalogue.json"
    if cat_path.exists():
        for e in json.loads(cat_path.read_text(encoding="utf-8")):
            src = e.get("source", "")
            raw_dir = None
            if src in ("HFWR", "HFWS") and e.get("number"):
                cand = SLIDES_DIR / src / e["number"]
                if cand.is_dir():
                    raw_dir = cand
            add({
                "name":   e["name"],
                "number": e.get("number", ""),
                "source": src,
                "pptx":   ROOT / e["pptx"] if e.get("pptx") else None,
                "raw_dir": raw_dir,
            })

    # 2) eChoice Output — folder per song, whole-deck only.
    ech = SLIDES_DIR / "eChoice Output"
    if ech.is_dir():
        for d in ech.iterdir():
            pptx = d / f"{d.name}.pptx"
            if pptx.exists():
                add({"name": d.name, "number": "", "source": "eChoice",
                     "pptx": pptx, "raw_dir": None})

    # 3) Hand-made top-level decks (hymn-slides/*.pptx) — whole-deck only.
    for pptx in SLIDES_DIR.glob("*.pptx"):
        add({"name": pptx.stem, "number": "", "source": "Custom",
             "pptx": pptx, "raw_dir": None})

    return index


def resolve_song(item: dict, index: dict) -> dict | None:
    """Find the best deck entry for a lineup song item."""
    title = item.get("title") or item.get("name", "")
    matches = index.get(normalize(title), [])
    if not matches:
        return None
    # Disambiguate by explicit source / number when given.
    want_src = (item.get("source") or "").upper()
    want_num = str(item.get("number") or "").lstrip("0")
    if want_src:
        narrowed = [m for m in matches if m["source"].upper() == want_src]
        matches = narrowed or matches
    if want_num:
        narrowed = [m for m in matches if m["number"].lstrip("0") == want_num]
        matches = narrowed or matches
    # Prefer entries that allow verse selection if verses were requested.
    if item.get("verses") not in (None, "all"):
        matches = sorted(matches, key=lambda m: m["raw_dir"] is None)
    return matches[0]


# ── Slide harvesting ────────────────────────────────────────────────────────────

def slide_sort_key(stem: str):
    """Order PNG slides: title, verses (n.m), chorus, then anything else."""
    s = stem.upper()
    if s.endswith("TITLE"):
        return (0, 0, 0)
    m = re.search(r"VERSE\s+(\d+)\.(\d+)", s)
    if m:
        return (1, int(m.group(1)), int(m.group(2)))
    m = re.search(r"CHORUS[.\s]+(\d+)", s)
    if m:
        return (2, 0, int(m.group(1)))
    m = re.search(r"(CODA|BRIDGE|REFRAIN|END)[.\s]*(\d*)", s)
    if m:
        return (3, 0, int(m.group(2) or 0))
    m = re.search(r"\.(\d+)$", s)
    if m:
        return (4, 0, int(m.group(1)))
    return (9, 0, 0)


def parse_verses(spec) -> set | None:
    """Turn a 'verses' value into a set of ints, or None for 'all'."""
    if spec in (None, "all", "ALL", ""):
        return None
    if isinstance(spec, int):
        return {spec}
    if isinstance(spec, (list, tuple)):
        out = set()
        for v in spec:
            out |= (parse_verses(v) or set())
        return out
    if isinstance(spec, str):
        out = set()
        for part in re.split(r"[,\s]+", spec.strip()):
            if not part:
                continue
            if "-" in part:                       # e.g. "1-3"
                a, b = part.split("-", 1)
                out |= set(range(int(a), int(b) + 1))
            else:
                out.add(int(part))
        return out or None
    return None


def pngs_for_song(raw_dir: Path, verses: set | None) -> list[Path]:
    """Selected PNG slides from an HFWR/HFWS source folder, in singing order."""
    pngs = [p for p in raw_dir.glob("*.png") if "PRINT" not in p.stem.upper()]
    pngs.sort(key=lambda p: slide_sort_key(p.stem))
    if verses is None:
        return pngs
    kept = []
    for p in pngs:
        s = p.stem.upper()
        m = re.search(r"VERSE\s+(\d+)\.", s)
        if m:
            if int(m.group(1)) in verses:
                kept.append(p)
        else:
            kept.append(p)        # title / chorus / coda / refrain always kept
    return kept


def pictures_from_pptx(pptx_path: Path):
    """Yield (ext, blob) for every full-bleed picture slide in a deck, in order."""
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        pics = [sh for sh in slide.shapes if sh.shape_type == PICTURE_SHAPE]
        if pics:
            img = pics[0].image
            yield img.ext, img.blob


# ── Slide construction ─────────────────────────────────────────────────────────

def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_full_picture(prs, source):
    """source: a Path to an image file or a (ext, blob) tuple."""
    slide = add_blank(prs)
    if isinstance(source, tuple):
        stream = io.BytesIO(source[1])
        slide.shapes.add_picture(stream, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    else:
        slide.shapes.add_picture(str(source), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    return slide


def _solid_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _textbox(slide, top, height, lines, font, size, color, *,
             bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(0.8), top, SLIDE_W - Inches(1.6), height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
    return box


def render_frame(prs, frame_key: str, template: dict, override: dict | None = None):
    style  = template["style"]
    frames = template["frames"]
    spec   = dict(frames.get(frame_key, {"title": frame_key.replace("_", " ").title()}))
    if override:
        spec.update({k: v for k, v in override.items()
                     if k in ("title", "subtitle", "lines", "background")})

    slide = add_blank(prs)
    bg_img = spec.get("background")
    if bg_img:
        path = (ROOT / bg_img)
        if path.exists():
            slide.shapes.add_picture(str(path), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
            return slide                          # artwork is self-contained
    _solid_bg(slide, style["bg"])

    title    = spec.get("title", "")
    subtitle = spec.get("subtitle", "")
    lines    = spec.get("lines", [])

    if lines:
        if title:
            _textbox(slide, Inches(0.7), Inches(1.4), [title],
                     style["title_font"], style["title_size"], style["accent"], bold=True)
        _textbox(slide, Inches(2.2), Inches(4.0), lines,
                 style["body_font"], style["body_size"], style["text"])
    else:
        _textbox(slide, Inches(2.6), Inches(2.3), [title] if title else [""],
                 style["title_font"], style["title_size"], style["accent"], bold=True)
        if subtitle:
            _textbox(slide, Inches(4.6), Inches(1.2), [subtitle],
                     style["body_font"], style["subtitle_size"], style["muted"])
    return slide


def render_scripture(prs, item: dict, template: dict):
    style = template["style"]
    title = item.get("title") or item.get("heading") or "Scripture Reading"
    ref   = item.get("reference", "")
    text  = item.get("text", "")
    # Allow text to be a list of verse strings or one block.
    if isinstance(text, list):
        body = list(text)
    else:
        body = [p for p in re.split(r"\n\s*\n", str(text).strip()) if p] or [""]

    # Title slide for the passage.
    s0 = add_blank(prs)
    _solid_bg(s0, style["bg"])
    _textbox(s0, Inches(2.6), Inches(2.3), [title],
             style["title_font"], style["title_size"], style["accent"], bold=True)
    if ref:
        _textbox(s0, Inches(4.7), Inches(1.0), [ref],
                 style["body_font"], style["subtitle_size"], style["muted"])

    # One slide per paragraph/verse so nothing overflows.
    for para in body:
        s = add_blank(prs)
        _solid_bg(s, style["bg"])
        _textbox(s, Inches(0.8), SLIDE_H - Inches(2.2), [para],
                 style["body_font"], style["body_size"], style["text"])
        if ref:
            _textbox(s, SLIDE_H - Inches(1.1), Inches(0.7), [ref],
                     style["body_font"], style["ref_size"], style["muted"])
    return 1 + len(body)


# ── Build ────────────────────────────────────────────────────────────────────

def build(lineup: dict, out_path: Path, template: dict, index: dict) -> dict:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    report = {"items": [], "warnings": [], "total_slides": 0}

    for i, item in enumerate(lineup.get("items", []), 1):
        itype = item.get("type", "song")
        label = item.get("title") or item.get("name") or item.get("frame") or itype
        added = 0

        if itype == "song":
            entry = resolve_song(item, index)
            if entry is None:
                report["warnings"].append(f"#{i}: song not found — '{label}'")
                report["items"].append({"n": i, "type": "song", "label": label,
                                        "status": "MISSING", "slides": 0})
                continue
            verses = parse_verses(item.get("verses"))
            if verses is not None and entry["raw_dir"] is None:
                report["warnings"].append(
                    f"#{i}: '{entry['name']}' ({entry['source']}) has no per-verse "
                    f"source — including the whole song.")
                verses = None
            if verses is not None and entry["raw_dir"] is not None:
                for png in pngs_for_song(entry["raw_dir"], verses):
                    add_full_picture(prs, png); added += 1
            else:
                if not entry["pptx"] or not Path(entry["pptx"]).exists():
                    report["warnings"].append(f"#{i}: deck file missing for '{label}'")
                    continue
                for src in pictures_from_pptx(Path(entry["pptx"])):
                    add_full_picture(prs, src); added += 1
            report["items"].append({
                "n": i, "type": "song", "label": entry["name"],
                "source": entry["source"],
                "verses": sorted(verses) if verses else "all",
                "slides": added})

        elif itype == "scripture":
            added = render_scripture(prs, item, template)
            report["items"].append({"n": i, "type": "scripture",
                                    "label": label, "slides": added})

        elif itype == "frame":
            key = item.get("frame") or item.get("name") or "blank"
            render_frame(prs, key, template, override=item)
            added = 1
            report["items"].append({"n": i, "type": "frame",
                                    "label": key, "slides": added})

        else:
            report["warnings"].append(f"#{i}: unknown item type '{itype}' — skipped")
            continue

        report["total_slides"] += added

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    report["out"] = str(out_path)
    return report


# ── CLI ────────────────────────────────────────────────────────────────────────

def default_out_name(lineup: dict, lineup_path: Path) -> Path:
    date = lineup.get("date", "")
    svc  = lineup.get("service", "CB Service")
    stem = f"{svc} ({date})".strip() if date else lineup_path.stem
    stem = re.sub(r'[<>:"/\\|?*]', "", stem)
    return OUT_DIR / f"{stem}.pptx"


def main():
    for _s in (sys.stdout, sys.stderr):       # survive the Windows cp1252 console
        try:
            _s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass
    ap = argparse.ArgumentParser(description="Build a Castleberry worship service deck.")
    ap.add_argument("lineup", nargs="?", help="Path to a lineup JSON spec.")
    ap.add_argument("-o", "--out", help="Output .pptx path.")
    ap.add_argument("--list-songs", action="store_true",
                    help="List every resolvable song title and exit.")
    args = ap.parse_args()

    index = build_song_index()

    if args.list_songs:
        names = sorted({m["name"] for ms in index.values() for m in ms})
        print(f"{len(names)} resolvable songs:\n")
        for n in names:
            print(" ", n)
        return

    if not args.lineup:
        ap.error("a lineup JSON path is required (or use --list-songs)")

    lineup_path = Path(args.lineup)
    lineup = json.loads(lineup_path.read_text(encoding="utf-8"))
    template = json.loads(TEMPLATE_PATH.read_text(encoding="utf-8"))

    out_path = Path(args.out) if args.out else default_out_name(lineup, lineup_path)
    report = build(lineup, out_path, template, index)

    print(f"\nBuilt: {report['out']}")
    print(f"Total slides: {report['total_slides']}\n")
    print("Order:")
    for it in report["items"]:
        extra = ""
        if it["type"] == "song":
            extra = f"  [{it.get('source','')}, verses={it.get('verses')}]"
        flag = "  ⚠ MISSING" if it.get("status") == "MISSING" else ""
        print(f"  {it['n']:>2}. {it['type']:<9} {it['label']}{extra} — {it['slides']} slides{flag}")
    if report["warnings"]:
        print("\nWarnings:")
        for w in report["warnings"]:
            print("  •", w)


if __name__ == "__main__":
    main()
